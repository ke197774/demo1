#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Supabase 技术分享 PPT 生成器
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 定义颜色主题
PRIMARY_COLOR = RgbColor(59, 130, 246)     # 蓝色
ACCENT_COLOR = RgbColor(16, 185, 129)      # 绿色
DARK_COLOR = RgbColor(30, 41, 59)          # 深灰
LIGHT_COLOR = RgbColor(248, 250, 252)      # 浅灰
SUPABASE_GREEN = RgbColor(61, 213, 152)    # Supabase 绿

def add_title_slide(title, subtitle=""):
    """添加标题幻灯片"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 背景形状
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_COLOR
    shape.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = SUPABASE_GREEN
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(title):
    """添加章节分隔幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 左侧装饰条
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = SUPABASE_GREEN
    shape.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = DARK_COLOR
    p.alignment = PP_ALIGN.LEFT

    return slide

def add_content_slide(title, bullets, note=""):
    """添加内容幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_COLOR

    # 分隔线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(2), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = SUPABASE_GREEN
    line.line.fill.background()

    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(bullet, tuple):
            # 带层级的子弹点
            p.text = bullet[0]
            p.level = bullet[1]
        else:
            p.text = f"• {bullet}"
            p.level = 0

        p.font.size = Pt(22)
        p.font.color.rgb = DARK_COLOR
        p.space_before = Pt(12)

    return slide

def add_feature_slide(title, features):
    """添加特性展示幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_COLOR

    # 特性卡片
    card_width = Inches(3.8)
    card_height = Inches(2.2)
    start_x = Inches(0.5)
    start_y = Inches(1.5)
    gap = Inches(0.4)

    colors = [SUPABASE_GREEN, PRIMARY_COLOR, ACCENT_COLOR, RgbColor(239, 68, 68)]

    for i, (feature_name, feature_desc) in enumerate(features):
        row = i // 3
        col = i % 3

        x = start_x + col * (card_width + gap)
        y = start_y + row * (card_height + gap)

        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RgbColor(248, 250, 252)
        card.line.color.rgb = colors[i % len(colors)]
        card.line.width = Pt(2)

        # 特性名称
        name_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), card_width - Inches(0.4), Inches(0.5))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = feature_name
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = colors[i % len(colors)]

        # 特性描述
        desc_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.7), card_width - Inches(0.4), card_height - Inches(0.9))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = feature_desc
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_COLOR

    return slide

def add_comparison_slide(title, headers, rows):
    """添加对比表格幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_COLOR

    # 表格
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.5), Inches(1.3), Inches(12.333), Inches(4)).table

    # 设置列宽
    col_widths = [Inches(2.5), Inches(3.277), Inches(3.277), Inches(3.277)]
    for i, width in enumerate(col_widths):
        table.columns[i].width = width

    # 表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = SUPABASE_GREEN
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = RgbColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # 表格内容
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RgbColor(248, 250, 252)

    return slide

def add_demo_slide(title, steps):
    """添加Demo步骤幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_COLOR

    # 步骤
    start_y = Inches(1.5)
    for i, (step_title, step_time, step_desc) in enumerate(steps):
        # 步骤编号圆圈
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.5), start_y + Inches(i * 1.1), Inches(0.6), Inches(0.6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = SUPABASE_GREEN
        circle.line.fill.background()

        # 编号文字
        num_box = slide.shapes.add_textbox(Inches(0.5), start_y + Inches(i * 1.1) + Inches(0.1), Inches(0.6), Inches(0.4))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RgbColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # 步骤标题
        step_box = slide.shapes.add_textbox(Inches(1.3), start_y + Inches(i * 1.1), Inches(8), Inches(0.5))
        tf = step_box.text_frame
        p = tf.paragraphs[0]
        p.text = step_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = DARK_COLOR

        # 时间标签
        time_box = slide.shapes.add_textbox(Inches(10), start_y + Inches(i * 1.1), Inches(2), Inches(0.5))
        tf = time_box.text_frame
        p = tf.paragraphs[0]
        p.text = step_time
        p.font.size = Pt(14)
        p.font.color.rgb = PRIMARY_COLOR
        p.alignment = PP_ALIGN.RIGHT

        # 步骤描述
        desc_box = slide.shapes.add_textbox(Inches(1.3), start_y + Inches(i * 1.1) + Inches(0.4), Inches(10.5), Inches(0.6))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = step_desc
        p.font.size = Pt(16)
        p.font.color.rgb = RgbColor(100, 116, 139)

    return slide

def add_code_slide(title, code_text):
    """添加代码展示幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_COLOR

    # 代码背景
    code_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5)
    )
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = RgbColor(30, 41, 59)
    code_bg.line.fill.background()

    # 代码文本
    code_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.8), Inches(5))
    tf = code_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(14)
    p.font.name = "Consolas"
    p.font.color.rgb = RgbColor(226, 232, 240)

    return slide

def add_summary_slide(key_points, resources):
    """添加总结幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "总结与资源"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_COLOR

    # 分隔线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(2), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = SUPABASE_GREEN
    line.line.fill.background()

    # 关键要点
    points_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(3))
    tf = points_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "关键要点"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_COLOR

    for point in key_points:
        p = tf.add_paragraph()
        p.text = f"✓ {point}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_COLOR
        p.space_before = Pt(12)

    # 资源链接
    res_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.8), Inches(4))
    tf = res_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "推荐资源"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_COLOR

    for resource in resources:
        p = tf.add_paragraph()
        p.text = f"→ {resource}"
        p.font.size = Pt(16)
        p.font.color.rgb = PRIMARY_COLOR
        p.space_before = Pt(12)

    return slide

# ============================================================
# 创建演示文稿内容
# ============================================================

# Slide 1: 封面
add_title_slide(
    "Supabase 技术分享",
    "开源的 Firebase 替代方案 | 让前端开发者也能做全栈"
)

# Slide 2: 目录
add_content_slide(
    "分享目录",
    [
        "问题引入：传统全栈开发的痛点",
        "BaaS 概念与 Supabase 定位",
        "Supabase 核心功能详解",
        "实战 Demo：10分钟搭建应用",
        "技术选型与对比分析",
        "总结与资源推荐"
    ]
)

# Slide 3: 问题引入
add_content_slide(
    "传统全栈开发的痛点",
    [
        "后端搭建复杂：服务器配置、数据库设计、API 开发...",
        "运维成本高：部署、监控、扩展、安全...",
        "认证系统繁琐：注册登录、权限管理、安全防护",
        ("前端开发者面临的困境", 1),
        ("想快速验证想法，却被后端劝退", 2),
        ("MVP 开发周期长，错过市场窗口", 2),
    ]
)

# Slide 4: BaaS 概念
add_content_slide(
    "BaaS：Backend as a Service",
    [
        "什么是 BaaS？",
        ("云端后端服务，开箱即用", 1),
        ("数据库、认证、存储、API 一站式提供", 1),
        "Firebase 的成功与局限",
        ("成功：快速开发、实时同步、Google 生态", 1),
        ("局限：闭源、NoSQL 限制、Vendor Lock-in", 1),
        "Supabase 的定位",
        ("开源的 Firebase 替代方案", 1),
        ("基于 PostgreSQL，功能强大且灵活", 1),
    ]
)

# Slide 5: Supabase 简介
add_content_slide(
    "Supabase 是什么？",
    [
        "开源的 Backend-as-a-Service 平台",
        "基于 PostgreSQL 构建的完整后端解决方案",
        "99.2K+ GitHub Stars，Top 100 开源项目",
        "被 Mozilla、GitHub、1Password 等企业使用",
        "SOC2 Type 2 认证，HIPAA 合规",
        "支持 React、Vue、Flutter、Next.js 等主流框架"
    ]
)

# Slide 6: 核心功能概览
add_feature_slide(
    "Supabase 核心功能",
    [
        ("PostgreSQL 数据库", "全功能 Postgres\n自动生成 REST API\n实时数据订阅"),
        ("认证系统 Auth", "邮箱/社交登录\nRow Level Security\nJWT Token 管理"),
        ("存储 Storage", "文件上传管理\nCDN 加速分发\n细粒度权限控制"),
        ("Edge Functions", "Deno 运行时\n全球边缘部署\nWebhook 支持"),
        ("实时订阅 Realtime", "WebSocket 连接\n多人协作场景\n变更监听"),
        ("向量数据库 Vector", "AI 向量存储\nEmbedding 搜索\nOpenAI 集成"),
    ]
)

# Slide 7: 数据库功能详解
add_content_slide(
    "PostgreSQL 数据库",
    [
        "每个项目 = 一个完整的 PostgreSQL 数据库",
        "自动生成 RESTful API",
        ("无需手写 CRUD 接口", 1),
        ("支持复杂查询、事务、存储过程", 1),
        "实时订阅（Realtime）",
        ("基于 PostgreSQL 的 LISTEN/NOTIFY", 1),
        ("WebSocket 连接，毫秒级同步", 1),
        "表编辑器：可视化建表、修改数据",
        "SQL 编辑器：支持复杂查询和调试"
    ]
)

# Slide 8: 认证系统
add_content_slide(
    "认证系统 (Authentication)",
    [
        "多种登录方式",
        ("邮箱 + 密码", 1),
        ("社交登录：Google、GitHub、Apple、微信...", 1),
        ("Magic Link：无密码登录", 1),
        ("手机号 + 短信验证码", 1),
        "Row Level Security (RLS)",
        ("行级权限控制，数据安全有保障", 1),
        ("声明式策略，SQL 级别实现", 1),
        "用户管理：Dashboard 可视化管理用户"
    ]
)

# Slide 9: 存储与 Functions
add_content_slide(
    "存储 & Edge Functions",
    [
        "Storage 文件存储",
        ("支持图片、视频等大文件", 1),
        ("自动 CDN 加速", 1),
        ("图片变换：缩放、裁剪、格式转换", 1),
        ("细粒度权限控制", 1),
        "Edge Functions 边缘函数",
        ("Deno 运行时，TypeScript 原生支持", 1),
        ("全球 300+ 节点部署", 1),
        ("适合 Webhook、定时任务、API 扩展", 1),
    ]
)

# Slide 10: Demo 介绍
add_demo_slide(
    "实战 Demo：待办事项应用",
    [
        ("创建 Supabase 项目", "1 分钟", "注册账号，创建新项目，获取 API 密钥"),
        ("创建数据表 + 配置 RLS", "2 分钟", "设计表结构，设置行级安全策略"),
        ("连接前端代码", "3 分钟", "安装 SDK，配置连接，实现 CRUD"),
        ("添加认证功能", "2 分钟", "集成登录组件，实现用户隔离"),
        ("展示实时同步", "2 分钟", "打开多窗口，演示实时数据同步效果"),
    ]
)

# Slide 11: 代码示例 - 初始化
add_code_slide(
    "代码示例：初始化项目",
    """// 安装依赖
npm install @supabase/supabase-js

// 初始化客户端
import { createClient } from '@supabase/supabase-js'

const supabaseUrl = 'https://your-project.supabase.co'
const supabaseKey = 'your-anon-key'

export const supabase = createClient(supabaseUrl, supabaseKey)

// 类型定义（TypeScript）
interface Todo {
  id: number
  title: string
  completed: boolean
  user_id: string
  created_at: string
}"""
)

# Slide 12: 代码示例 - CRUD
add_code_slide(
    "代码示例：增删改查",
    """// 查询数据
const { data, error } = await supabase
  .from('todos')
  .select('*')
  .eq('user_id', userId)
  .order('created_at', { ascending: false })

// 插入数据
const { data, error } = await supabase
  .from('todos')
  .insert({ title: 'Learn Supabase', user_id: userId })

// 更新数据
const { error } = await supabase
  .from('todos')
  .update({ completed: true })
  .eq('id', todoId)

// 删除数据
const { error } = await supabase
  .from('todos')
  .delete()
  .eq('id', todoId)"""
)

# Slide 13: 代码示例 - 实时订阅
add_code_slide(
    "代码示例：实时订阅",
    """// 订阅数据变化
const channel = supabase
  .channel('todos-changes')
  .on(
    'postgres_changes',
    {
      event: '*',          // INSERT, UPDATE, DELETE
      schema: 'public',
      table: 'todos',
      filter: `user_id=eq.${userId}`
    },
    (payload) => {
      console.log('数据变化:', payload)
      // 实时更新 UI
      if (payload.eventType === 'INSERT') {
        addTodoToList(payload.new)
      } else if (payload.eventType === 'DELETE') {
        removeTodoFromList(payload.old)
      }
    }
  )
  .subscribe()

// 取消订阅
channel.unsubscribe()"""
)

# Slide 14: 代码示例 - 认证
add_code_slide(
    "代码示例：用户认证",
    """// 邮箱注册
const { data, error } = await supabase.auth.signUp({
  email: 'user@example.com',
  password: 'password123'
})

// 邮箱登录
const { data, error } = await supabase.auth.signInWithPassword({
  email: 'user@example.com',
  password: 'password123'
})

// GitHub 社交登录
const { data, error } = await supabase.auth.signInWithOAuth({
  provider: 'github'
})

// 获取当前用户
const { data: { user } } = await supabase.auth.getUser()

// 登出
await supabase.auth.signOut()"""
)

# Slide 15: 技术对比
add_comparison_slide(
    "技术选型对比",
    ["对比项", "Supabase", "Firebase", "自建后端"],
    [
        ["开源", "✅ 完全开源", "❌ 闭源", "✅ 自主可控"],
        ["数据库", "PostgreSQL (SQL)", "Firestore (NoSQL)", "任意选择"],
        ["学习成本", "低", "低", "高"],
        ["迁移成本", "低（可导出）", "中（有锁定）", "低"],
        ["自托管", "✅ 支持", "❌ 不支持", "✅ 支持"],
        ["适用场景", "中小型项目", "快速原型", "大型系统"],
    ]
)

# Slide 16: 选型建议
add_content_slide(
    "选型建议",
    [
        "推荐使用 Supabase 的场景",
        ("个人项目、创业 MVP", 1),
        ("前端团队独立完成全栈", 1),
        ("需要 SQL 数据库复杂查询", 1),
        ("关注数据隐私，希望可自托管", 1),
        "其他选择",
        ("Firebase：已有 Firebase 项目，或需要 NoSQL", 1),
        ("自建后端：复杂业务逻辑、高并发场景", 1),
    ]
)

# Slide 17: 价格
add_content_slide(
    "价格方案",
    [
        "Free 免费版",
        ("500MB 数据库存储", 1),
        ("1GB 文件存储", 1),
        ("50,000 月活用户", 1),
        ("无限 API 请求", 1),
        "Pro 专业版 $25/月",
        ("8GB 数据库存储", 1),
        ("100GB 文件存储", 1),
        ("100,000 月活用户", 1),
        ("每日备份 + 优先支持", 1),
    ]
)

# Slide 18: 总结
add_summary_slide(
    [
        "Supabase = 开源 Firebase 替代方案",
        "基于 PostgreSQL，功能强大灵活",
        "前后端一站式，开发效率提升 10x",
        "免费额度友好，适合快速验证想法",
        "开源可控，可自托管无锁定风险"
    ],
    [
        "官网: supabase.com",
        "文档: supabase.com/docs",
        "GitHub: github.com/supabase/supabase",
        "Discord 社区: supabase.com/discord",
        "示例项目: github.com/supabase/supabase/tree/master/examples"
    ]
)

# Slide 19: Q&A
add_title_slide(
    "Q & A",
    "感谢聆听 | 欢迎提问"
)

# 保存演示文稿
output_path = "d:/tmp/teach/Supabase技术分享.pptx"
prs.save(output_path)
print(f"PPT 已生成: {output_path}")
print(f"共 {len(prs.slides)} 张幻灯片")
